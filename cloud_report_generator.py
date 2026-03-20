from reportlab.lib.pagesizes import letter
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer,
    Table, TableStyle, PageBreak, Image
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
import tempfile
from datetime import datetime, timezone
import pandas as pd

# PowerPoint support
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from io import BytesIO

print("Report generator loaded")


def generate_boardroom_pdf(client, monthly_spend=None, savings_monthly=None, top_service_name=None, maturity_score=None, readiness_score=None):
    """
    Generate a professional boardroom-style PDF with just client parameter.
    
    Args:
        client (str): Client name
    
    Returns:
        str: Path to the generated PDF file
    """
    # Default values for demo purposes
    monthly_spend = 150000 if monthly_spend is None else monthly_spend
    savings_monthly = 25000 if savings_monthly is None else savings_monthly
    top_service_name = "EC2 Instances" if top_service_name is None else top_service_name
    maturity_score = 78 if maturity_score is None else maturity_score
    readiness_score = 85 if readiness_score is None else readiness_score
    
    file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(file.name, pagesize=letter)
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        "Title",
        fontSize=22,
        leading=28,
        alignment=1,
        spaceAfter=20
    )

    section_style = ParagraphStyle(
        "Section",
        fontSize=16,
        spaceAfter=12,
        textColor=colors.darkblue
    )

    content = []

    # ===== COVER PAGE =====
    content.append(Spacer(1, 120))
    content.append(Paragraph("Cloud Executive Advisory Report", title_style))
    content.append(Spacer(1, 40))

    content.append(Paragraph(f"Client: {client}", styles["Normal"]))
    content.append(Paragraph(f"Date: {datetime.now(timezone.utc).date()}", styles["Normal"]))
    content.append(Paragraph("Confidential — For Executive Review", styles["Italic"]))

    # --- Executive KPI Cards (KPI style) ---
    from reportlab.platypus import Flowable
    class CardRow(Flowable):
        def __init__(self, metrics, width=500, height=70, gap=18):
            Flowable.__init__(self)
            self.metrics = metrics
            self.width = width
            self.height = height
            self.gap = gap
            self.card_colors = [
                colors.HexColor("#0066cc"),
                colors.HexColor("#00994c"),
                colors.HexColor("#ff9900"),
                colors.HexColor("#6600cc")
            ]
        def draw(self):
            c = self.canv
            n = len(self.metrics)
            card_w = (self.width - (n-1)*self.gap) / n
            x = 0
            for i, (title, value, subtitle) in enumerate(self.metrics):
                c.saveState()
                c.setFillColor(self.card_colors[i % len(self.card_colors)])
                c.roundRect(x, 0, card_w, self.height, 14, fill=1, stroke=0)
                c.setFillColor(colors.white)
                c.setFont("Helvetica-Bold", 13)
                c.drawCentredString(x + card_w/2, self.height-22, title)
                c.setFont("Helvetica-Bold", 20)
                c.drawCentredString(x + card_w/2, self.height-44, value)
                if subtitle:
                    c.setFont("Helvetica", 11)
                    c.drawCentredString(x + card_w/2, 12, subtitle)
                c.restoreState()
                x += card_w + self.gap
        def wrap(self, availWidth, availHeight):
            return (self.width, self.height)
    percent_savings = (savings_monthly/monthly_spend*100) if monthly_spend else 0
    metrics = [
        ("Monthly Cloud Spend", f"₹{monthly_spend:,.0f}", ""),
        ("Estimated Savings", f"₹{savings_monthly:,.0f} ({percent_savings:.1f}%)", ""),
        ("Cloud Maturity", f"{maturity_score}/100", ""),
        ("Transformation Readiness", f"{readiness_score}/100", ""),
        ("Focus", "Cost optimization and modernization", "")
    ]
    content.append(Spacer(1, 30))
    content.append(CardRow(metrics, width=500, height=70, gap=18))
    content.append(Spacer(1, 24))
    content.append(PageBreak())

    # ===== PRIORITY ACTIONS =====
    content.append(Paragraph("Priority Actions — Next 30 Days", section_style))

    actions = [
        "Optimize high-cost service footprint",
        "Implement Savings Plans / Reserved Instances",
        "Remove idle resources and unused assets"
    ]

    for action in actions:
        content.append(Paragraph(f"• {action}", styles["Normal"]))

    content.append(PageBreak())

    # ===== FINANCIAL IMPACT =====
    content.append(Paragraph("Financial Impact", section_style))

    content.append(Paragraph("Cost Distribution", section_style))
    content.append(Image("cost_distribution.png", width=400, height=300))

    content.append(Paragraph(
        f"<b>Monthly Spend:</b> ₹{monthly_spend:,.0f}",
        ParagraphStyle("KPI", fontSize=16, spaceAfter=12)
    ))

    content.append(Paragraph(
        f"Estimated Annual Savings: ₹{savings_monthly*12:,.0f}",
        ParagraphStyle("Highlight",
            backColor=colors.lightblue,
            borderPadding=8,
            fontSize=12
        )
    ))
    content.append(Paragraph(
        "Estimated ROI Timeline: Less than 12 months",
        styles["Normal"]
    ))
    content.append(Paragraph(
        "Savings can be reinvested into innovation and modernization initiatives.",
        styles["Normal"]
    ))

    content.append(PageBreak())

    content.append(Paragraph("Executive Dashboard — " + client, section_style))

    content.append(Paragraph(
        f"Cloud Maturity: {maturity_score}/100",
        styles["Normal"]
    ))

    content.append(Paragraph(
        f"Transformation Readiness: {readiness_score}/100",
        styles["Normal"]
    ))

    doc.build(content)

    return file.name


def generate_executive_pdf(client):
    """
    Generate a simplified executive PDF with just client parameter.
    
    Args:
        client (str): Client name
    
    Returns:
        str: Path to the generated PDF file
    """
    # Default values for demo purposes
    monthly_spend = 150000
    savings_monthly = 25000
    top_service_name = "EC2 Instances"
    maturity_score = 78
    
    file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(file.name, pagesize=letter)
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        "Title",
        fontSize=22,
        leading=28,
        alignment=1,
        spaceAfter=20
    )

    section_style = ParagraphStyle(
        "Section",
        fontSize=16,
        spaceAfter=12,
        textColor=colors.darkblue
    )

    content = []

    # ===== COVER PAGE =====
    content.append(Spacer(1, 120))
    content.append(Paragraph("Cloud Executive Advisory Report", title_style))
    content.append(Spacer(1, 40))

    content.append(Paragraph(f"Client: {client}", styles["Normal"]))
    content.append(Paragraph(f"Date: {datetime.now(timezone.utc).date()}", styles["Normal"]))
    content.append(Paragraph("Confidential — For Executive Review", styles["Italic"]))

    content.append(PageBreak())

    # ===== EXECUTIVE SNAPSHOT =====
    content.append(Paragraph("Executive Snapshot", ParagraphStyle(
        "Header",
        fontSize=18,
        textColor=colors.darkblue,
        spaceAfter=12
    )))

    snapshot_data = [
        ["Metric", "Value"],
        ["Monthly Spend", f"₹{monthly_spend:,.0f}"],
        ["Annual Run Rate", f"₹{monthly_spend*12:,.0f}"],
        ["Savings Opportunity", f"₹{savings_monthly:,.0f}"],
        ["Top Cost Driver", top_service_name],
        ["Cloud Maturity", f"{maturity_score}/100"],
        ["Transformation Readiness", f"{readiness_score}/100"]
    ]

    snapshot_table = Table(snapshot_data, colWidths=[220, 180])

    snapshot_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.darkblue),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1),
            [colors.whitesmoke, colors.lightgrey]),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey)
    ]))

    content.append(snapshot_table)
    content.append(PageBreak())

    # ===== PRIORITY ACTIONS =====
    content.append(Paragraph("Priority Actions — Next 30 Days", section_style))

    actions = [
        "Optimize high-cost service footprint",
        "Implement Savings Plans / Reserved Instances",
        "Remove idle resources and unused assets"
    ]

    for action in actions:
        content.append(Paragraph(f"• {action}", styles["Normal"]))

    content.append(PageBreak())

    # ===== FINANCIAL IMPACT =====
    content.append(Paragraph("Financial Impact", section_style))

    content.append(Paragraph("Cost Distribution", section_style))
    content.append(Image("cost_distribution.png", width=400, height=300))

    content.append(Paragraph(
        f"<b>Monthly Spend:</b> ₹{monthly_spend:,.0f}",
        ParagraphStyle("KPI", fontSize=16, spaceAfter=12)
    ))

    content.append(Paragraph(
        f"Estimated Annual Savings: ₹{savings_monthly*12:,.0f}",
        ParagraphStyle("Highlight",
            backColor=colors.lightblue,
            borderPadding=8,
            fontSize=12
        )
    ))
    content.append(Paragraph(
        "Estimated ROI Timeline: Less than 12 months",
        styles["Normal"]
    ))
    content.append(Paragraph(
        "Savings can be reinvested into innovation and modernization initiatives.",
        styles["Normal"]
    ))

    content.append(PageBreak())

    content.append(Paragraph("Executive Dashboard — " + client, section_style))

    content.append(Paragraph(
        f"Cloud Maturity: {maturity_score}/100",
        styles["Normal"]
    ))

    content.append(Paragraph(
        f"Transformation Readiness: {readiness_score}/100",
        styles["Normal"]
    ))

    doc.build(content)

    return file.name


def generate_boardroom_pdf(client, monthly_spend=None, savings_monthly=None, top_service_name=None, maturity_score=None, readiness_score=None):
    """
    Generate a professional boardroom-style PDF cloud executive report.
    
    Args:
        client (str): Client name
        monthly_spend (float, optional): Monthly cloud spend amount
        savings_monthly (float, optional): Potential monthly savings
        top_service_name (str, optional): Name of top cost-driving service
        maturity_score (int, optional): Cloud maturity score (0-100)
    
    Returns:
        str: Path to the generated PDF file
    """
    # Use default values if not provided
    if monthly_spend is None:
        monthly_spend = 150000
    if savings_monthly is None:
        savings_monthly = 25000
    if top_service_name is None:
        top_service_name = "EC2 Instances"
    if maturity_score is None:
        maturity_score = 78
    if readiness_score is None:
        readiness_score = 85
    
    file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(file.name, pagesize=letter)
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        "Title",
        fontSize=22,
        leading=28,
        alignment=1,
        spaceAfter=20
    )

    section_style = ParagraphStyle(
        "Section",
        fontSize=16,
        spaceAfter=12,
        textColor=colors.darkblue
    )

    content = []

    # ===== COVER PAGE =====
    content.append(Spacer(1, 120))
    content.append(Paragraph("Cloud Executive Advisory Report", title_style))
    content.append(Spacer(1, 40))

    content.append(Paragraph(f"Client: {client}", styles["Normal"]))
    content.append(Paragraph(f"Date: {datetime.now(timezone.utc).date()}", styles["Normal"]))
    content.append(Paragraph("Confidential — For Executive Review", styles["Italic"]))

    content.append(PageBreak())

    # ===== EXECUTIVE SNAPSHOT =====
    content.append(Paragraph("Executive Snapshot", ParagraphStyle(
        "Header",
        fontSize=18,
        textColor=colors.darkblue,
        spaceAfter=12
    )))

    snapshot_data = [
        ["Metric", "Value"],
        ["Monthly Spend", f"₹{monthly_spend:,.0f}"],
        ["Annual Run Rate", f"₹{monthly_spend*12:,.0f}"],
        ["Savings Opportunity", f"₹{savings_monthly:,.0f}"],
        ["Top Cost Driver", top_service_name],
        ["Cloud Maturity", f"{maturity_score}/100"]
    ]

    snapshot_table = Table(snapshot_data, colWidths=[220, 180])

    snapshot_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.darkblue),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 11),
        ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1),
            [colors.whitesmoke, colors.lightgrey]),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey)
    ]))

    content.append(snapshot_table)
    content.append(PageBreak())

    # ===== PRIORITY ACTIONS =====
    content.append(Paragraph("Priority Actions — Next 30 Days", section_style))

    actions = [
        "Optimize high-cost service footprint",
        "Implement Savings Plans / Reserved Instances",
        "Remove idle resources and unused assets"
    ]

    for action in actions:
        content.append(Paragraph(f"• {action}", styles["Normal"]))

    content.append(PageBreak())

    # ===== FINANCIAL IMPACT =====
    content.append(Paragraph("Financial Impact", section_style))

    content.append(Paragraph("Cost Distribution", section_style))
    content.append(Image("cost_distribution.png", width=400, height=300))

    content.append(Paragraph(
        f"<b>Monthly Spend:</b> ₹{monthly_spend:,.0f}",
        ParagraphStyle("KPI", fontSize=16, spaceAfter=12)
    ))

    content.append(Paragraph(
        f"Estimated Annual Savings: ₹{savings_monthly*12:,.0f}",
        ParagraphStyle("Highlight",
            backColor=colors.lightblue,
            borderPadding=8,
            fontSize=12
        )
    ))
    content.append(Paragraph(
        "Estimated ROI Timeline: Less than 12 months",
        styles["Normal"]
    ))
    content.append(Paragraph(
        "Savings can be reinvested into innovation and modernization initiatives.",
        styles["Normal"]
    ))

    content.append(PageBreak())

    content.append(Paragraph("Executive Dashboard — " + client, section_style))

    content.append(Paragraph(
        f"Cloud Maturity: {maturity_score}/100",
        styles["Normal"]
    ))

    content.append(Paragraph(
        f"Transformation Readiness: {readiness_score}/100",
        styles["Normal"]
    ))

    doc.build(content)

    return file.name

def generate_excel_report(client, monthly_spend=None, savings_monthly=None, top_service_name=None, maturity_score=None, service_cost=None, df=None):
    """
    Generate a comprehensive Excel report with multiple sheets.
    
    Args:
        client (str): Client name
        monthly_spend (float, optional): Monthly cloud spend amount
        savings_monthly (float, optional): Potential monthly savings
        top_service_name (str, optional): Name of top cost-driving service
        maturity_score (int, optional): Cloud maturity score (0-100)
        service_cost (pd.DataFrame, optional): Service cost breakdown
        df (pd.DataFrame, optional): Raw billing data
    
    Returns:
        str: Path to the generated Excel file
    """
    # Use default values if not provided
    if monthly_spend is None:
        monthly_spend = 150000
    if savings_monthly is None:
        savings_monthly = 25000
    if top_service_name is None:
        top_service_name = "EC2 Instances"
    if maturity_score is None:
        maturity_score = 78
    
    file = tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx")

    with pd.ExcelWriter(file.name, engine="openpyxl") as writer:

        # ===== Sheet 1 — Executive Summary =====
        summary_df = pd.DataFrame({
            "Metric": [
                "Client",
                "Date",
                "Monthly Spend",
                "Annual Run Rate",
                "Savings Opportunity",
                "Top Cost Driver",
                "Cloud Maturity"
            ],
            "Value": [
                client,
                datetime.now(timezone.utc).date(),
                f"₹{monthly_spend:,.0f}",
                f"₹{monthly_spend*12:,.0f}",
                f"₹{savings_monthly:,.0f}",
                top_service_name,
                f"{maturity_score}/100"
            ]
        })

        summary_df.to_excel(writer, sheet_name="Executive Summary", index=False)
        
        # Format Executive Summary columns
        worksheet = writer.book["Executive Summary"]
        worksheet.column_dimensions["A"].width = 30
        worksheet.column_dimensions["B"].width = 25

        # ===== Sheet 2 — Top Services =====
        if service_cost is not None:
            service_cost.to_excel(writer, sheet_name="Cost by Service", index=False)
        else:
            # Create sample service cost data
            sample_services = pd.DataFrame({
                "Service": ["EC2 Instances", "RDS Database", "S3 Storage", "CloudFront", "Lambda"],
                "Cost": [75000, 30000, 22500, 15000, 7500]
            })
            sample_services.to_excel(writer, sheet_name="Cost by Service", index=False)

        # ===== Sheet 3 — Raw Data =====
        if df is not None:
            df.to_excel(writer, sheet_name="Raw Billing Data", index=False)
        else:
            # Create sample raw data
            sample_raw = pd.DataFrame({
                "Service": ["EC2", "RDS", "S3", "CloudFront", "Lambda"],
                "Cost": [75000, 30000, 22500, 15000, 7500],
                "Usage": ["High", "Medium", "Low", "Medium", "Low"]
            })
            sample_raw.to_excel(writer, sheet_name="Raw Billing Data", index=False)

    return file.name

def generate_powerpoint_report(
    client,
    monthly_spend=None,
    savings_monthly=None,
    top_service_name=None,
    maturity_score=None,
    readiness_score=85
):
    """
    Generate a modern executive PowerPoint deck with charts and KPIs.
    Args:
        client (str): Client name
        monthly_spend (float, optional): Monthly cloud spend amount
        savings_monthly (float, optional): Potential monthly savings
        top_service_name (str, optional): Name of top cost-driving service
        maturity_score (int, optional): Cloud maturity score (0-100)
        readiness_score (int, optional): Transformation readiness score (0-100)
        service_cost (pd.DataFrame, optional): Service cost breakdown
        trend_data (pd.DataFrame, optional): Monthly trend data (date, cost)
    Returns:
        str: Path to the generated PowerPoint file
    """
    import matplotlib.pyplot as plt
    import tempfile
    import os
    if monthly_spend is None:
        monthly_spend = 150000
    if savings_monthly is None:
        savings_monthly = 25000
    if top_service_name is None:
        top_service_name = "EC2 Instances"
    if maturity_score is None:
        maturity_score = 78
    readiness_score = 85 if readiness_score is None else readiness_score

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_slide_layout)

    # --- Slide 1: Executive Summary ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Executive Summary"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Client: {client}    |    Date: {datetime.now(timezone.utc).strftime('%Y-%m-%d')}"
    subtitle_frame.paragraphs[0].font.size = Pt(16)
    subtitle_frame.paragraphs[0].font.italic = True
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Metric Cards (centered)

    metrics = [
        ("Monthly Cloud Spend", f"${monthly_spend:,.0f}", RGBColor(0, 102, 204)),
        ("Estimated Savings", f"${savings_monthly:,.0f} ({savings_monthly/monthly_spend*100:.1f}%)", RGBColor(0, 153, 76)),
        ("Cloud Maturity", f"{maturity_score}/100", RGBColor(255, 153, 0)),
        ("Focus", "Cost optimization and modernization", RGBColor(102, 0, 204)),
    ]
    card_width = Inches(5.5)
    card_height = Inches(1.1)
    y_gap = Inches(0.35)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    total_height = len(metrics) * card_height + (len(metrics) - 1) * y_gap
    x_left = (slide_width - card_width) / 2
    y_top = (slide_height - total_height) / 2 + Inches(0.3)  # slight downward offset for title
    for i, (label, value, color) in enumerate(metrics):
        card = slide.shapes.add_shape(
            5,  # Rounded Rectangle
            x_left,
            y_top + i * (card_height + y_gap),
            card_width,
            card_height
        )
        fill = card.fill
        fill.solid()
        fill.fore_color.rgb = color
        card.line.color.rgb = RGBColor(255, 255, 255)
        card.line.width = Pt(2)
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        # Center text vertically in card
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.LEFT

    # --- Slide 2: Cost Distribution (Bar and Pie Chart) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf2 = title2.text_frame
    tf2.text = "Cost Distribution"
    tf2.paragraphs[0].font.size = Pt(34)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER


    # Track temp image files for later cleanup
    temp_img_files = []

    # Bar chart (service_cost)
    bar_img = None
    if service_cost is not None and not service_cost.empty:
        fig, ax = plt.subplots(figsize=(6, 3))
        ax.barh(service_cost['Service'], service_cost['Cost'], color='#007acc')
        ax.set_xlabel('Cost')
        ax.set_title('Cost by Service')
        plt.tight_layout()
        bar_imgfile = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        plt.savefig(bar_imgfile.name, bbox_inches='tight')
        plt.close(fig)
        bar_img = bar_imgfile.name
        slide2.shapes.add_picture(bar_img, Inches(1.0), Inches(1.2), Inches(6.5), Inches(2.2))
        temp_img_files.append(bar_img)

    # Pie chart (service_cost)
    if service_cost is not None and not service_cost.empty:
        fig, ax = plt.subplots(figsize=(3, 3))
        ax.pie(service_cost['Cost'], labels=service_cost['Service'], autopct='%1.0f%%', startangle=140)
        ax.set_title('Service Mix')
        plt.tight_layout()
        pie_imgfile = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        plt.savefig(pie_imgfile.name, bbox_inches='tight')
        plt.close(fig)
        slide2.shapes.add_picture(pie_imgfile.name, Inches(7.2), Inches(1.2), Inches(2.5), Inches(2.2))
        temp_img_files.append(pie_imgfile.name)

    # --- Slide 3: Key Performance Indicators (Grid) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf3 = title3.text_frame
    tf3.text = "Key Performance Indicators"
    tf3.paragraphs[0].font.size = Pt(34)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER

    kpis = [
        ("Monthly Spend", f"${monthly_spend:,.0f}", RGBColor(0, 102, 204)),
        ("Annual Run Rate", f"${monthly_spend*12:,.0f}", RGBColor(0, 153, 76)),
        ("Savings Opportunity", f"${savings_monthly:,.0f}", RGBColor(255, 153, 0)),
        ("Top Cost Driver", top_service_name, RGBColor(102, 0, 204)),
        ("Cloud Maturity", f"{maturity_score}/100", RGBColor(204, 0, 0)),
    ]
    grid_x = [1.0, 5.0]
    grid_y = [1.5, 2.7, 3.9]
    for idx, (label, value, color) in enumerate(kpis):
        x = grid_x[idx % 2]
        y = grid_y[idx // 2]
        card = slide3.shapes.add_shape(
            5,  # Rounded Rectangle
            Inches(x),
            Inches(y),
            Inches(3.8),
            Inches(1.0)
        )
        fill = card.fill
        fill.solid()
        fill.fore_color.rgb = color
        card.line.color.rgb = RGBColor(255, 255, 255)
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.LEFT

    # --- Slide 4: Executive Dashboard ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf4 = title4.text_frame
    tf4.text = "Executive Dashboard"
    tf4.paragraphs[0].font.size = Pt(34)
    tf4.paragraphs[0].font.bold = True
    tf4.paragraphs[0].alignment = PP_ALIGN.CENTER

    dash_metrics = [
        ("Cloud Maturity", f"{maturity_score}/100", RGBColor(0, 102, 204)),
        ("Transformation Readiness", f"{readiness_score}/100", RGBColor(0, 153, 76)),
        ("Optimization Priority", "High", RGBColor(255, 153, 0)),
    ]
    for i, (label, value, color) in enumerate(dash_metrics):
        card = slide4.shapes.add_shape(
            5,  # Rounded Rectangle
            Inches(1.0 + i * 3.1),
            Inches(1.5),
            Inches(3.0),
            Inches(1.2)
        )
        fill = card.fill
        fill.solid()
        fill.fore_color.rgb = color
        card.line.color.rgb = RGBColor(255, 255, 255)
        card.line.width = Pt(1.5)
        tf = card.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.LEFT

    # --- Slide 5: Trends & Opportunities (optional) ---

    if trend_data is not None and not trend_data.empty:
        slide5 = prs.slides.add_slide(prs.slide_layouts[6])
        title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf5 = title5.text_frame
        tf5.text = "Trends & Opportunities"
        tf5.paragraphs[0].font.size = Pt(34)
        tf5.paragraphs[0].font.bold = True
        tf5.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Line chart
        fig, ax = plt.subplots(figsize=(6, 3))
        ax.plot(trend_data['Date'], trend_data['Cost'], marker='o', color='#007acc')
        ax.set_xlabel('Date')
        ax.set_ylabel('Cost')
        ax.set_title('Monthly Spend Trend')
        plt.xticks(rotation=30)
        plt.tight_layout()
        trend_imgfile = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        plt.savefig(trend_imgfile.name, bbox_inches='tight')
        plt.close(fig)
        slide5.shapes.add_picture(trend_imgfile.name, Inches(1.0), Inches(1.2), Inches(7.5), Inches(2.2))
        temp_img_files.append(trend_imgfile.name)


    # Save to temp file
    file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(file.name)

    # Now cleanup temp images
    for img_path in temp_img_files:
        try:
            os.unlink(img_path)
        except Exception:
            pass

    return file.name
    """
    Generate a comprehensive Excel report with multiple sheets.
    
    Args:
        client (str): Client name
        monthly_spend (float, optional): Monthly cloud spend amount
        savings_monthly (float, optional): Potential monthly savings
        top_service_name (str, optional): Name of top cost-driving service
        maturity_score (int, optional): Cloud maturity score (0-100)
        service_cost (pd.DataFrame, optional): Service cost breakdown
        df (pd.DataFrame, optional): Raw billing data
    
    Returns:
        str: Path to the generated Excel file
    """
    # Use default values if not provided
    if monthly_spend is None:
        monthly_spend = 150000
    if savings_monthly is None:
        savings_monthly = 25000
    if top_service_name is None:
        top_service_name = "EC2 Instances"
    if maturity_score is None:
        maturity_score = 78
    
    file = tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx")

    with pd.ExcelWriter(file.name, engine="openpyxl") as writer:

        # ===== Sheet 1 — Executive Summary =====
        summary_df = pd.DataFrame({
            "Metric": [
                "Client",
                "Date",
                "Monthly Spend",
                "Annual Run Rate",
                "Savings Opportunity",
                "Top Cost Driver",
                "Cloud Maturity"
            ],
            "Value": [
                client,
                datetime.now(timezone.utc).strftime('%Y-%m-%d'),
                f"${monthly_spend:,.0f}",
                f"${monthly_spend*12:,.0f}",
                f"${savings_monthly:,.0f}",
                top_service_name,
                f"{maturity_score}/100"
            ]
        })

        summary_df.to_excel(writer, sheet_name="Executive Summary", index=False)
        worksheet = writer.book["Executive Summary"]
        worksheet.column_dimensions["A"].width = 30
        worksheet.column_dimensions["B"].width = 25
        # Left-align all cells in Executive Summary
        for row in worksheet.iter_rows(min_row=2, max_row=worksheet.max_row, min_col=2, max_col=2):
            for cell in row:
                cell.alignment = cell.alignment.copy(horizontal='left')

        # ===== Sheet 2 — Top Services =====
        if service_cost is not None:
            service_cost_fmt = service_cost.copy()
            if 'Cost' in service_cost_fmt.columns:
                service_cost_fmt['Cost'] = service_cost_fmt['Cost'].apply(lambda x: f"${x:,.2f}")
            service_cost_fmt.to_excel(writer, sheet_name="Cost by Service", index=False)
        else:
            sample_services = pd.DataFrame({
                "Service": ["EC2 Instances", "RDS Database", "S3 Storage", "CloudFront", "Lambda"],
                "Cost": [f"${v:,.2f}" for v in [75000, 30000, 22500, 15000, 7500]]
            })
            sample_services.to_excel(writer, sheet_name="Cost by Service", index=False)

        # ===== Sheet 3 — Raw Data =====
        if df is not None:
            df_fmt = df.copy()
            if 'cost' in df_fmt.columns:
                df_fmt['cost'] = df_fmt['cost'].apply(lambda x: f"${x:,.2f}")
            df_fmt.to_excel(writer, sheet_name="Raw Billing Data", index=False)
        else:
            sample_raw = pd.DataFrame({
                "Service": ["EC2", "RDS", "S3", "CloudFront", "Lambda"],
                "Cost": [f"${v:,.2f}" for v in [75000, 30000, 22500, 15000, 7500]],
                "Usage": ["High", "Medium", "Low", "Medium", "Low"]
            })
            sample_raw.to_excel(writer, sheet_name="Raw Billing Data", index=False)

    return file.name


def generate_dashboard_pdf(client, monthly_spend=None, savings_monthly=None, maturity_score=None, readiness_score=None):
    """
    Generate a dashboard-style PDF report with charts and metrics.
    
    Args:
        client (str): Client name
        monthly_spend (float, optional): Monthly cloud spend amount
        savings_monthly (float, optional): Potential monthly savings
        maturity_score (int, optional): Cloud maturity score (0-100)
        readiness_score (int, optional): Transformation readiness score (0-100)
    
    Returns:
        str: Path to the generated PDF file
    """
    # Use default values if not provided
    if monthly_spend is None:
        monthly_spend = 150000
    if savings_monthly is None:
        savings_monthly = 25000
    if maturity_score is None:
        maturity_score = 78
    if readiness_score is None:
        readiness_score = 70

    file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(file.name, pagesize=letter)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "Title", fontSize=20, alignment=1, spaceAfter=20
    )

    section_style = ParagraphStyle(
        "Section", fontSize=16, textColor=colors.darkblue, spaceAfter=10
    )

    content = []

    # ===== COVER =====
    content.append(Spacer(1, 120))
    content.append(Paragraph("Cloud Executive Dashboard Report", title_style))
    content.append(Paragraph(f"Client: {client}", styles["Normal"]))
    content.append(Paragraph(f"Date: {datetime.now(timezone.utc).date()}", styles["Normal"]))
    content.append(PageBreak())

    # ===== EXECUTIVE DASHBOARD =====
    content.append(Paragraph(f"Executive Dashboard — {client}", section_style))

    content.append(Paragraph(
        f"Monthly Spend: ₹{monthly_spend:,.0f}", styles["Normal"]
    ))
    content.append(Paragraph(
        f"Savings Opportunity: ₹{savings_monthly:,.0f}", styles["Normal"]
    ))
    content.append(Paragraph(
        f"Cloud Maturity: {maturity_score}/100", styles["Normal"]
    ))
    content.append(Paragraph(
        f"Transformation Readiness: {readiness_score}/100", styles["Normal"]
    ))

    content.append(PageBreak())

    # ===== COST DISTRIBUTION CHART =====
    content.append(Paragraph("Cost Distribution", section_style))
    content.append(Image("cost_distribution.png", width=400, height=300))

    content.append(PageBreak())

    # ===== COST BY SERVICE =====
    content.append(Paragraph("Cost by Service", section_style))
    content.append(Image("cost_by_service.png", width=400, height=300))

    content.append(PageBreak())

    # ===== PRIORITY ACTIONS =====
    content.append(Paragraph("Priority Actions", section_style))

    actions = [
        "Optimize high-cost services",
        "Implement Savings Plans",
        "Remove idle resources"
    ]

    for action in actions:
        content.append(Paragraph(f"• {action}", styles["Normal"]))

    doc.build(content)

    return file.name


# Legacy function for backward compatibility
def generate_pdf(monthly_spend, savings_monthly, top_service_name, maturity_score, top_services_data=None):
    """
    Generate a PDF cloud executive report (legacy function).
    
    Args:
        monthly_spend (float): Monthly cloud spend amount
        savings_monthly (float): Potential monthly savings
        top_service_name (str): Name of top cost-driving service
        maturity_score (int): Cloud maturity score (0-100)
        top_services_data (list, optional): List of tuples with service data for table
    
    Returns:
        str: Path to the generated PDF file
    """
    file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")

    doc = SimpleDocTemplate(file.name, pagesize=letter)
    styles = getSampleStyleSheet()

    content = []

    content.append(Paragraph("Cloud Executive Report", styles["Title"]))
    content.append(Spacer(1, 12))

    content.append(Paragraph(f"Monthly Spend: ₹{monthly_spend:,.0f}", styles["Normal"]))
    content.append(Paragraph(f"Savings Opportunity: ₹{savings_monthly:,.0f}", styles["Normal"]))
    content.append(Paragraph(f"Top Cost Driver: {top_service_name}", styles["Normal"]))
    content.append(Paragraph(f"Cloud Maturity Score: {maturity_score}/100", styles["Normal"]))
    
    if top_services_data:
        content.append(Spacer(1, 12))
        content.append(Paragraph("Top Services", styles["Heading2"]))
        
        # Create table using tabulate
        table_text = tabulate(top_services_data, headers=["Service", "Cost (₹)", "Percentage"], 
                              tablefmt="grid", floatfmt=".2f")
        
        # Add table as paragraph (simple approach)
        content.append(Paragraph(table_text, styles["Code"]))

    content.append(PageBreak())

    content.append(Paragraph("Executive Dashboard — " + client, section_style))

    content.append(Paragraph(
        f"Cloud Maturity: {maturity_score}/100",
        styles["Normal"]
    ))

    content.append(Paragraph(
        f"Transformation Readiness: {readiness_score}/100",
        styles["Normal"]
    ))

    doc.build(content)

    return file.name

# Example usage
if __name__ == "__main__":
    # Sample data
    monthly_spend = 150000
    savings_monthly = 25000
    top_service_name = "EC2 Instances"
    maturity_score = 78
    client = "Acme Corp"
    
    # Generate executive PDF
    pdf_path = generate_executive_pdf(client, monthly_spend, savings_monthly, top_service_name, maturity_score)
    print(f"Executive PDF generated: {pdf_path}")

# Streamlit integration example
def streamlit_export_button(client, monthly_spend, savings_monthly, top_service_name, maturity_score, top_services_data=None):
    """
    Streamlit sidebar button and download functionality for PDF export.
    
    Args:
        client (str): Client name
        monthly_spend (float): Monthly cloud spend amount
        savings_monthly (float): Potential monthly savings
        top_service_name (str): Name of top cost-driving service
        maturity_score (int): Cloud maturity score (0-100)
        top_services_data (list, optional): List of tuples with service data for table
    """
    import streamlit as st

    if st.sidebar.button("Export Executive Report (PDF)"):
        pdf_file = generate_executive_pdf(client, monthly_spend, savings_monthly, top_service_name, maturity_score)
        with open(pdf_file, "rb") as f:
            st.sidebar.download_button(
                label="Download Report",
                data=f,
                file_name="Cloud_Executive_Report.pdf",
                mime="application/pdf"
            )

    if st.sidebar.button("Export Executive Report (PowerPoint)"):
        pptx_file = generate_powerpoint_report(client, monthly_spend, savings_monthly, top_service_name, maturity_score)
        with open(pptx_file, "rb") as f:
            st.sidebar.download_button(
                label="Download PowerPoint",
                data=f,
                file_name="Cloud_Executive_Report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

# Example usage
if __name__ == "__main__":
    # Sample data
    monthly_spend = 150000
    savings_monthly = 25000
    top_service_name = "EC2 Instances"
    maturity_score = 78
    
    # Sample top services data for table
    top_services_data = [
        ("EC2 Instances", 75000, 50.0),
        ("RDS Database", 30000, 20.0),
        ("S3 Storage", 22500, 15.0),
        ("CloudFront", 15000, 10.0),
        ("Lambda", 7500, 5.0)
    ]
    
    # Generate PDF
    pdf_path = generate_pdf(monthly_spend, savings_monthly, top_service_name, maturity_score, top_services_data)
    print(f"PDF generated: {pdf_path}")

# Streamlit integration example
def streamlit_export_button(monthly_spend, savings_monthly, top_service_name, maturity_score, top_services_data=None):
    """
    Streamlit sidebar button and download functionality for PDF export.
    
    Args:
        monthly_spend (float): Monthly cloud spend amount
        savings_monthly (float): Potential monthly savings
        top_service_name (str): Name of top cost-driving service
        maturity_score (int): Cloud maturity score (0-100)
        top_services_data (list, optional): List of tuples with service data for table
    """
    import streamlit as st
    
    if st.sidebar.button("Export Executive Report (PDF)"):
        pdf_file = generate_pdf(monthly_spend, savings_monthly, top_service_name, maturity_score, top_services_data)

        with open(pdf_file, "rb") as f:
            st.sidebar.download_button(
                label="Download Report",
                data=f,
                file_name="Cloud_Executive_Report.pdf",
                mime="application/pdf"
            )
