from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
from datetime import datetime, timezone

def create_colored_shape(slide, left, top, width, height, color_rgb, text=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color_rgb)
    shape.line.fill.background()
    if text:
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    return shape

def create_dashboard_image(monthly_spend, savings_monthly, maturity_score, readiness_score):
    plt.figure(figsize=(10,5))
    plt.text(0.1, 0.8, f"Monthly Spend\n₹{monthly_spend:,.0f}", fontsize=16)
    plt.text(0.5, 0.8, f"Savings\n₹{savings_monthly:,.0f}", fontsize=16)
    plt.text(0.1, 0.4, f"Maturity\n{maturity_score}/100", fontsize=16)
    plt.text(0.5, 0.4, f"Readiness\n{readiness_score}/100", fontsize=16)
    plt.axis('off')
    plt.savefig("executive_dashboard.png")
    plt.close()
    return "executive_dashboard.png"

def generate_executive_ppt(client, monthly_spend=None, savings_monthly=None, maturity_score=None, readiness_score=None, service_cost=None):
    """
    Generate an executive PowerPoint presentation with charts and metrics.
    
    Args:
        client (str): Client name
        monthly_spend (float, optional): Monthly cloud spend amount
        savings_monthly (float, optional): Potential monthly savings
        maturity_score (int, optional): Cloud maturity score (0-100)
        readiness_score (int, optional): Transformation readiness score (0-100)
        service_cost (pd.DataFrame, optional): Service cost breakdown
    
    Returns:
        str: Path to the generated PowerPoint file
    """
    # --- LOCKED SECTION: Executive Snapshot (Slide 2) ---
    # This section is locked as per user request. Do not modify without explicit approval.
    def _locked_slide_modification():
        raise RuntimeError("Executive Snapshot (slide 2) is locked and cannot be modified without explicit approval.")

    # Use default values if not provided
    if monthly_spend is None:
        monthly_spend = 150000
    if savings_monthly is None:
        savings_monthly = 25000
    if maturity_score is None:
        maturity_score = 78
    if readiness_score is None:
        readiness_score = 70

    # Debug: print/log service_cost DataFrame
    try:
        import logging
        logging.basicConfig(level=logging.INFO)
        logging.info("[Executive PPTX] service_cost DataFrame:\n%s", str(service_cost))
        print("[Executive PPTX] service_cost DataFrame:\n", service_cost)
    except Exception as e:
        print("[Executive PPTX] Could not print service_cost DataFrame:", e)

    # Create presentation
    prs = Presentation()

    # ===== SLIDE 1: CUSTOM DESIGN TITLE SLIDE (REFINED) =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout for full control

    # Main dark blue background
    bg_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(11, 25, 44)  # #0b192c
    bg_shape.line.fill.background()

    # (No accent circles; removed as per user request)

    # Executive Summary label (top left, smaller)
    label_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(4), Inches(0.3))
    label_frame = label_box.text_frame
    label_frame.clear()
    p = label_frame.paragraphs[0]
    p.text = "EXECUTIVE SUMMARY"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.name = "Montserrat"
    p.font.color.rgb = RGBColor(26, 188, 156)
    p.space_after = 0

    # Main Title (top left, smaller)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(6), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "Cloud Executive\nAdvisory Report"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Montserrat"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = 0

    # Client label (top left, below title)
    client_label = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(1.2), Inches(0.2))
    client_label_frame = client_label.text_frame
    client_label_frame.clear()
    p = client_label_frame.paragraphs[0]
    p.text = "Client"
    p.font.size = Pt(12)
    p.font.name = "Open Sans"
    p.font.color.rgb = RGBColor(143, 160, 181)
    p.space_after = 0

    # Client name (top left, below client label)
    client_name_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(2.5), Inches(0.3))
    client_name_frame = client_name_box.text_frame
    client_name_frame.clear()
    p = client_name_frame.paragraphs[0]
    p.text = str(client)
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = "Montserrat"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = 0

    # Date label (top left, right of client)
    date_label = slide.shapes.add_textbox(Inches(3.2), Inches(2.2), Inches(1.2), Inches(0.2))
    date_label_frame = date_label.text_frame
    date_label_frame.clear()
    p = date_label_frame.paragraphs[0]
    p.text = "Date"
    p.font.size = Pt(12)
    p.font.name = "Open Sans"
    p.font.color.rgb = RGBColor(143, 160, 181)
    p.space_after = 0

    # Date value (top left, right of client name)
    date_value_box = slide.shapes.add_textbox(Inches(3.2), Inches(2.4), Inches(2), Inches(0.3))
    date_value_frame = date_value_box.text_frame
    date_value_frame.clear()
    p = date_value_frame.paragraphs[0]
    p.text = str(datetime.now(timezone.utc).date())
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = "Montserrat"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = 0

    # Bottom bar (dark blue)
    bottom_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0, prs.slide_height - Inches(0.7), prs.slide_width, Inches(0.7)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = RGBColor(8, 18, 32)  # #081220
    bottom_bar.line.fill.background()

    # Confidential text
    conf_box = slide.shapes.add_textbox(0, prs.slide_height - Inches(0.45), prs.slide_width, Inches(0.3))
    conf_frame = conf_box.text_frame
    conf_frame.clear()
    p = conf_frame.paragraphs[0]
    p.text = "Confidential — For Executive Review Only"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = "Open Sans"
    p.font.color.rgb = RGBColor(231, 76, 60)  # #e74c3c
    p.alignment = PP_ALIGN.CENTER
    p.space_after = 0

    # Decorative dots (simulate with text, more subtle)
    dots_box = slide.shapes.add_textbox(Inches(11.5), Inches(1.2), Inches(1.2), Inches(1.2))
    dots_frame = dots_box.text_frame
    dots_frame.clear()
    p = dots_frame.paragraphs[0]
    p.text = "+ + + +\n+ + + +\n+ + + +\n+ + + +"
    p.font.size = Pt(18)
    p.font.name = "Montserrat"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    p.space_after = 0
    
    # ===== SLIDE 2: EXECUTIVE SNAPSHOT (CUSTOM DESIGN) =====
    # Slide 2: Executive Snapshot (locked)
    add_executive_snapshot_slide(prs, monthly_spend, savings_monthly, maturity_score, readiness_score, service_cost)
    # ... (rest of the slides and logic) ...
    # Save presentation
    file_path = f"{client.replace(' ', '_')}_executive_presentation.pptx"
    prs.save(file_path)
    return file_path

def add_executive_snapshot_slide(prs, monthly_spend, savings_monthly, maturity_score, readiness_score, service_cost):
    """
    Adds a custom 'Executive Snapshot' slide with KPI cards, matching the provided design.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(11, 25, 44)  # #0b192c
    bg_shape.line.fill.background()


    # Decorative dots (simulate with text)
    dots_box = slide.shapes.add_textbox(Inches(12.0), Inches(0.7), Inches(1.2), Inches(1.2))
    dots_frame = dots_box.text_frame
    dots_frame.clear()
    p = dots_frame.paragraphs[0]
    p.text = "+ + +\n+ + +\n+ + +"
    p.font.size = Pt(18)
    p.font.name = "Montserrat"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    p.space_after = 0

    # Header (left-aligned, closer to left edge)
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(8), Inches(0.6))
    header_frame = header_box.text_frame
    header_frame.clear()
    p = header_frame.paragraphs[0]
    p.text = "Executive Snapshot"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Montserrat"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    p.space_after = 0

    # Removed green (teal) underline below the Executive Snapshot title

    # KPI Card Data
    # Fix Top Cost Driver value logic to show actual value if available
    # Robust extraction of Top Cost Driver
    import numpy as np
    top_cost_driver = "-"
    if service_cost is not None and not service_cost.empty and "Service" in service_cost.columns:
        # Drop NaN and blank values, get the first valid service name
        valid_services = service_cost["Service"].dropna().astype(str)
        valid_services = valid_services[valid_services.str.strip() != ""]
        if not valid_services.empty:
            top_cost_driver = valid_services.iloc[0]
    kpis = [
        {
            "label": "MONTHLY SPEND",
            "value": f"${monthly_spend:,.0f}",
            "color": RGBColor(52, 152, 219),
            "icon": "wallet"
        },
        {
            "label": "ANNUAL RUN RATE",
            "value": f"${monthly_spend*12:,.0f}",
            "color": RGBColor(52, 152, 219),
            "icon": "calendar"
        },
        {
            "label": "SAVINGS OPPORTUNITY",
            "value": f"${savings_monthly:,.0f}",
            "color": RGBColor(26, 188, 156),
            "icon": "piggy"
        },
        {
            "label": "TOP COST DRIVER",
            "value": top_cost_driver,
            "color": RGBColor(52, 152, 219),
            "icon": "server"
        },
        {
            "label": "CLOUD MATURITY",
            "value": f"{maturity_score}/100",
            "color": RGBColor(26, 188, 156),
            "icon": "gauge"
        },
    ]


    # Card positions (x, y) - visually balanced, non-overlapping, matching design
    # Top row: 3 cards, evenly spaced
    # Bottom row: 2 cards, centered below
    card_positions = [
        (Inches(0.4), Inches(1.7)),   # Monthly Spend (left)
        (Inches(3.25), Inches(1.7)),  # Annual Run Rate (center)
        (Inches(6.1), Inches(1.7)),   # Savings Opportunity (right)
        (Inches(1.3), Inches(3.7)),   # Top Cost Driver (bottom left)
        (Inches(5.0), Inches(3.7)),   # Cloud Maturity (bottom right)
    ]

    card_width = Inches(2.7)
    card_height = Inches(1.7)


    for i, kpi in enumerate(kpis):
        x, y = card_positions[i]
        # Card background
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x, y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(21, 38, 63)  # #15263f
        card.shadow.inherit = False
        card.shadow.blur_radius = 20
        card.shadow.distance = 10
        card.shadow.transparency = 0.8
        card.line.fill.background()

        # Icon (centered)
        icon_box = slide.shapes.add_textbox(x + Inches(1.1), y + Inches(0.25), Inches(0.5), Inches(0.5))
        icon_frame = icon_box.text_frame
        icon_frame.clear()
        p = icon_frame.paragraphs[0]
        icon_map = {
            "wallet": "💼",
            "calendar": "📅",
            "piggy": "🐷",
            "server": "🖥️",
            "gauge": "⏱️"
        }
        p.text = icon_map.get(kpi["icon"], "")
        p.font.size = Pt(20)
        p.font.name = "Segoe UI Emoji"
        p.font.color.rgb = kpi["color"]
        p.space_after = 0
        p.alignment = PP_ALIGN.CENTER

        # Label (centered)
        label_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.85), Inches(2.3), Inches(0.3))
        label_frame = label_box.text_frame
        label_frame.clear()
        p = label_frame.paragraphs[0]
        p.text = kpi["label"]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = "Open Sans"
        p.font.color.rgb = kpi["color"] if kpi["label"] == "SAVINGS OPPORTUNITY" or kpi["label"] == "CLOUD MATURITY" else RGBColor(143, 160, 181)
        p.space_after = 0
        p.alignment = PP_ALIGN.CENTER

        # Value (centered)
        value_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.15), Inches(2.3), Inches(0.7))
        value_frame = value_box.text_frame
        value_frame.clear()
        p = value_frame.paragraphs[0]
        p.text = kpi["value"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.name = "Montserrat"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = 0
        p.alignment = PP_ALIGN.CENTER

        # Removed green progress bar under Cloud Maturity KPI card
    
    # ===== SLIDE 3: COST DISTRIBUTION CHART =====
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Cost Distribution"
    
    # Add pie chart image if it exists
    try:
        img_path = "cost_distribution.png"
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(5))
    except FileNotFoundError:
        # Fallback text if image not found
        content = slide.placeholders[1]
        content.text = "Cost distribution chart will be displayed here"
    
    # ===== SLIDE 4: KEY METRICS =====
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Key Performance Indicators"
    
    # Create metrics text
    metrics_text = f"""Financial Metrics:
• Monthly Spend: ₹{monthly_spend:,.0f}
• Annual Run Rate: ₹{monthly_spend*12:,.0f}
• Savings Opportunity: ₹{savings_monthly:,.0f}

Organizational Metrics:
• Cloud Maturity Score: {maturity_score}/100
• Transformation Readiness: {readiness_score}/100
• Optimization Priority: High"""
    
    content = slide.placeholders[1]
    content.text = metrics_text
    
    # ===== SLIDE 5: EXECUTIVE DASHBOARD =====
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Executive Dashboard"
    
    # Add dashboard image if it exists
    try:
        img_path = "executive_dashboard.png"
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(5))
    except FileNotFoundError:
        # Fallback text if image not found
        content = slide.placeholders[1]
        content.text = "Executive dashboard visualization will be displayed here"
    
    # ===== SLIDE 6: RECOMMENDATIONS =====
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Strategic Recommendations"
    
    recommendations_text = """Immediate Actions (0-30 Days):
• Identify and remove idle resources
• Implement cost monitoring alerts
• Rightsize over-provisioned instances

Short-term Initiatives (1-3 Months):
• Implement Reserved Instances/Savings Plans
• Optimize storage lifecycle policies
• Establish FinOps governance

Long-term Strategy (3-12 Months):
• Modernize legacy workloads
• Adopt cloud-native architectures
• Implement continuous optimization framework"""
    
    content = slide.placeholders[1]
    content.text = recommendations_text
    




def generate_partner_board_pack(
    client,
    monthly_spend,
    savings_monthly,
    maturity_score,
    readiness_score,
    top_service,
    service_cost
):
    """
    Generates a Partner-Level Board Pack PowerPoint
    """

    annual_spend = monthly_spend * 12
    annual_savings = savings_monthly * 12

    # ===== Create Cost Distribution Chart =====
    labels = service_cost["Service"]
    values = service_cost["Cost"]

    plt.figure(figsize=(6, 6))
    plt.pie(values, labels=labels, autopct="%1.1f%%")
    plt.title("Cloud Cost Distribution")

    chart_path = "cost_distribution.png"
    plt.savefig(chart_path, bbox_inches="tight")
    plt.close()

    # ===== Create Presentation =====
    prs = Presentation()

    # ---------- Slide 1: Cover ----------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Cloud Transformation Board Pack"
    subtitle = slide.placeholders[1]
    subtitle.text = f"{client}\nPartner Executive Briefing"

    # ---------- Slide 2: Executive Summary ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"

    content = slide.placeholders[1]
    content.text = (
        f"• Annual Cloud Spend: ₹{annual_spend:,.0f}\n"
        f"• Optimization Potential: ₹{annual_savings:,.0f} annually\n"
        f"• Primary Cost Driver: {top_service}\n"
        "• Immediate action can deliver ROI within 12 months"
    )

    # ---------- Slide 3: Executive Dashboard ----------
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_box.text_frame.text = "Executive Dashboard"

    # Create dashboard image using the dedicated function
    create_dashboard_image(monthly_spend, savings_monthly, maturity_score, readiness_score)

    kpis = [
        ("Monthly Spend", f"₹{monthly_spend:,.0f}"),
        ("Savings Potential", f"₹{savings_monthly:,.0f}"),
        ("Cloud Maturity", f"{maturity_score}/100"),
        ("Transformation Readiness", f"{readiness_score}/100"),
    ]

    left = 0.5
    for label, value in kpis:
        box = slide.shapes.add_textbox(
            Inches(left), Inches(1.2), Inches(2.2), Inches(1.2)
        )
        tf = box.text_frame
        tf.text = label
        p = tf.add_paragraph()
        p.text = value
        p.font.size = Pt(20)
        p.font.bold = True
        left += 2.3

    # ---------- Slide 4: Cost Insights ----------
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    ).text_frame.text = "Cost Insights"

    slide.shapes.add_picture(
        chart_path, Inches(1), Inches(1.2), height=Inches(4)
    )

    # ---------- Slide 5: Risk & Benchmark ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risk & Benchmark Assessment"

    content = slide.placeholders[1]
    content.text = (
        "• High dependency on key services increases financial risk\n"
        "• Spend concentration exceeds typical benchmarks\n"
        "• Governance improvements recommended"
    )

    # ---------- Slide 6: Priority Actions ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Priority Actions (Next 90 Days)"

    content = slide.placeholders[1]
    content.text = (
        "1. Optimize high-cost services\n"
        "2. Implement Savings Plans / Reserved Instances\n"
        "3. Remove idle resources\n"
        "4. Establish FinOps governance"
    )

    # ---------- Slide 7: Investment & ROI ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Investment & ROI Outlook"

    content = slide.placeholders[1]
    content.text = (
        f"Projected Annual Spend: ₹{annual_spend:,.0f}\n"
        f"Potential Annual Savings: ₹{annual_savings:,.0f}\n"
        "ROI Timeline: < 12 months\n"
        "Savings can fund modernization initiatives"
    )

    # ---------- Slide 8: Transformation Roadmap ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Transformation Roadmap"

    content = slide.placeholders[1]
    content.text = (
        "Phase 1: Stabilize & Optimize (0–6 months)\n"
        "Phase 2: Modernize & Automate (6–18 months)\n"
        "Phase 3: Innovate & Transform (18+ months)"
    )

    output_path = "Partner_Level_Board_Pack.pptx"
    prs.save(output_path)

    return output_path


print("PowerPoint report generator loaded")

