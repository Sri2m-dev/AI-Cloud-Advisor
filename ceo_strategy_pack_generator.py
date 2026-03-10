from pptx import Presentation


def generate_ceo_strategy_pack(
    client,
    monthly_spend,
    savings_monthly,
    maturity_score,
    readiness_score,
    top_service
):

    annual_spend = monthly_spend * 12
    annual_savings = savings_monthly * 12

    prs = Presentation()

    # ===== Slide 1 — Title =====
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Enterprise Cloud & Digital Strategy"
    slide.placeholders[1].text = f"{client}\nCEO Strategy Brief"

    # ===== Slide 2 — Executive Narrative =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Narrative"

    slide.placeholders[1].text = (
        "Current cloud expenditure presents a significant optimization opportunity "
        f"of ₹{annual_savings:,.0f} annually. Redirecting these savings can fund "
        "innovation initiatives, enhance customer experience, and improve operational resilience."
    )

    # ===== Slide 3 — Business Impact =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Business Impact"

    slide.placeholders[1].text = (
        "• Improved customer experience through reliable digital platforms\n"
        "• Increased operational efficiency via automation\n"
        "• Faster innovation cycles\n"
        "• Enhanced financial performance"
    )

    # ===== Slide 4 — Strategic Risk =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risk of Inaction"

    slide.placeholders[1].text = (
        "• Rising operational costs\n"
        "• Reduced competitiveness\n"
        "• Slower innovation\n"
        "• Increased dependency on legacy architecture"
    )

    # ===== Slide 5 — Value Creation Opportunity =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Value Creation Opportunity"

    slide.placeholders[1].text = (
        f"Annual Cloud Spend: ₹{annual_spend:,.0f}\n"
        f"Optimization Potential: ₹{annual_savings:,.0f}\n"
        "Reinvestment into innovation and modernization"
    )

    # ===== Slide 6 — Competitive Advantage =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Competitive Advantage"

    slide.placeholders[1].text = (
        "• Faster time to market\n"
        "• Scalable digital capabilities\n"
        "• Data-driven decision making\n"
        "• AI-enabled operations"
    )

    # ===== Slide 7 — Transformation Roadmap =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Transformation Roadmap"

    slide.placeholders[1].text = (
        "Phase 1: Stabilize & Optimize\n"
        "Phase 2: Modernize & Automate\n"
        "Phase 3: Innovate & Transform"
    )

    # ===== Slide 8 — Investment Case =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Investment Case"

    slide.placeholders[1].text = (
        f"Cloud Maturity: {maturity_score}/100\n"
        f"Transformation Readiness: {readiness_score}/100\n"
        "Investment required to unlock long-term value"
    )

    output = "CEO_Strategy_Pack.pptx"
    prs.save(output)

    return output
