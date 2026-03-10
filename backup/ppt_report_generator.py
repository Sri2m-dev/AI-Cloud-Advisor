from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
from datetime import date

def create_colored_shape(slide, left, top, width, height, color_rgb, text=None):
    """
    Create a colored rectangle shape on a PowerPoint slide.
    
    Args:
        slide: PowerPoint slide object
        left (float): Left position in inches
        top (float): Top position in inches
        width (float): Width in inches
        height (float): Height in inches
        color_rgb (tuple): RGB color tuple (R, G, B)
        text (str, optional): Text to add to shape
    
    Returns:
        shape: PowerPoint shape object
    """
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color_rgb)
    shape.line.fill.background()
    
    if text:
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    
    return shape


def create_dashboard_image(monthly_spend, savings_monthly, maturity_score, readiness_score):
    """
    Create executive dashboard image with key metrics.
    
    Args:
        monthly_spend (float): Monthly cloud spend amount
        savings_monthly (float): Potential monthly savings
        maturity_score (int): Cloud maturity score (0-100)
        readiness_score (int): Transformation readiness score (0-100)
    
    Returns:
        str: Path to generated image file
    """
    plt.figure(figsize=(10,5))
    plt.text(0.1, 0.8, f"Monthly Spend\n₹{monthly_spend:,.0f}", fontsize=16)
    plt.text(0.5, 0.8, f"Savings\n₹{savings_monthly:,.0f}", fontsize=16)
    plt.text(0.1, 0.4, f"Maturity\n{maturity_score}/100", fontsize=16)
    plt.text(0.5, 0.4, f"Readiness\n{readiness_score}/100", fontsize=16)
    plt.axis('off')
    plt.savefig("executive_dashboard.png")
    plt.close()
    
    return "executive_dashboard.png"

def generate_executive_ppt(client, monthly_spend=None, savings_monthly=None, maturity_score=None, readiness_score=None, service_cost=None):
    """
    Generate an executive PowerPoint presentation with charts and metrics.
    
    Args:
        client (str): Client name
        monthly_spend (float, optional): Monthly cloud spend amount
        savings_monthly (float, optional): Potential monthly savings
        maturity_score (int, optional): Cloud maturity score (0-100)
        readiness_score (int, optional): Transformation readiness score (0-100)
        service_cost (pd.DataFrame, optional): Service cost breakdown
    
    Returns:
        str: Path to the generated PowerPoint file
    """
    # Use default values if not provided
    if monthly_spend is None:
        monthly_spend = 150000
    if savings_monthly is None:
        savings_monthly = 25000
    if maturity_score is None:
        maturity_score = 78
    if readiness_score is None:
        readiness_score = 70
    
    # Create presentation
    prs = Presentation()
    
    # ===== SLIDE 1: TITLE SLIDE =====
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Cloud Executive Dashboard"
    subtitle.text = f"Client: {client}\nDate: {date.today()}"
    
    # ===== SLIDE 2: EXECUTIVE SUMMARY =====
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Executive Summary"
    
    # Add bullet points
    text_frame = content.text_frame
    text_frame.text = f"""• Monthly Cloud Spend: ₹{monthly_spend:,.0f}
• Estimated Savings: ₹{savings_monthly:,.0f} ({(savings_monthly/monthly_spend)*100:.1f}%)
• Cloud Maturity: {maturity_score}/100
• Transformation Readiness: {readiness_score}/100
• Focus: Cost optimization and modernization"""
    
    # ===== SLIDE 3: COST DISTRIBUTION CHART =====
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Cost Distribution"
    
    # Add pie chart image if it exists
    try:
        img_path = "cost_distribution.png"
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(5))
    except FileNotFoundError:
        # Fallback text if image not found
        content = slide.placeholders[1]
        content.text = "Cost distribution chart will be displayed here"
    
    # ===== SLIDE 4: KEY METRICS =====
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Key Performance Indicators"
    
    # Create metrics text
    metrics_text = f"""Financial Metrics:
• Monthly Spend: ₹{monthly_spend:,.0f}
• Annual Run Rate: ₹{monthly_spend*12:,.0f}
• Savings Opportunity: ₹{savings_monthly:,.0f}

Organizational Metrics:
• Cloud Maturity Score: {maturity_score}/100
• Transformation Readiness: {readiness_score}/100
• Optimization Priority: High"""
    
    content = slide.placeholders[1]
    content.text = metrics_text
    
    # ===== SLIDE 5: EXECUTIVE DASHBOARD =====
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Executive Dashboard"
    
    # Add dashboard image if it exists
    try:
        img_path = "executive_dashboard.png"
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(5))
    except FileNotFoundError:
        # Fallback text if image not found
        content = slide.placeholders[1]
        content.text = "Executive dashboard visualization will be displayed here"
    
    # ===== SLIDE 6: RECOMMENDATIONS =====
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Strategic Recommendations"
    
    recommendations_text = """Immediate Actions (0-30 Days):
• Identify and remove idle resources
• Implement cost monitoring alerts
• Rightsize over-provisioned instances

Short-term Initiatives (1-3 Months):
• Implement Reserved Instances/Savings Plans
• Optimize storage lifecycle policies
• Establish FinOps governance

Long-term Strategy (3-12 Months):
• Modernize legacy workloads
• Adopt cloud-native architectures
• Implement continuous optimization framework"""
    
    content = slide.placeholders[1]
    content.text = recommendations_text
    
    # Save presentation
    file_path = f"{client.replace(' ', '_')}_executive_presentation.pptx"
    prs.save(file_path)
    
    return file_path


def generate_partner_board_pack(
    client,
    monthly_spend,
    savings_monthly,
    maturity_score,
    readiness_score,
    top_service,
    service_cost
):
    """
    Generates a Partner-Level Board Pack PowerPoint
    """

    annual_spend = monthly_spend * 12
    annual_savings = savings_monthly * 12

    # ===== Create Cost Distribution Chart =====
    labels = service_cost["Service"]
    values = service_cost["Cost"]

    plt.figure(figsize=(6, 6))
    plt.pie(values, labels=labels, autopct="%1.1f%%")
    plt.title("Cloud Cost Distribution")

    chart_path = "cost_distribution.png"
    plt.savefig(chart_path, bbox_inches="tight")
    plt.close()

    # ===== Create Presentation =====
    prs = Presentation()

    # ---------- Slide 1: Cover ----------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Cloud Transformation Board Pack"
    subtitle = slide.placeholders[1]
    subtitle.text = f"{client}\nPartner Executive Briefing"

    # ---------- Slide 2: Executive Summary ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"

    content = slide.placeholders[1]
    content.text = (
        f"• Annual Cloud Spend: ₹{annual_spend:,.0f}\n"
        f"• Optimization Potential: ₹{annual_savings:,.0f} annually\n"
        f"• Primary Cost Driver: {top_service}\n"
        "• Immediate action can deliver ROI within 12 months"
    )

    # ---------- Slide 3: Executive Dashboard ----------
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_box.text_frame.text = "Executive Dashboard"

    # Create dashboard image using the dedicated function
    create_dashboard_image(monthly_spend, savings_monthly, maturity_score, readiness_score)

    kpis = [
        ("Monthly Spend", f"₹{monthly_spend:,.0f}"),
        ("Savings Potential", f"₹{savings_monthly:,.0f}"),
        ("Cloud Maturity", f"{maturity_score}/100"),
        ("Transformation Readiness", f"{readiness_score}/100"),
    ]

    left = 0.5
    for label, value in kpis:
        box = slide.shapes.add_textbox(
            Inches(left), Inches(1.2), Inches(2.2), Inches(1.2)
        )
        tf = box.text_frame
        tf.text = label
        p = tf.add_paragraph()
        p.text = value
        p.font.size = Pt(20)
        p.font.bold = True
        left += 2.3

    # ---------- Slide 4: Cost Insights ----------
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    ).text_frame.text = "Cost Insights"

    slide.shapes.add_picture(
        chart_path, Inches(1), Inches(1.2), height=Inches(4)
    )

    # ---------- Slide 5: Risk & Benchmark ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risk & Benchmark Assessment"

    content = slide.placeholders[1]
    content.text = (
        "• High dependency on key services increases financial risk\n"
        "• Spend concentration exceeds typical benchmarks\n"
        "• Governance improvements recommended"
    )

    # ---------- Slide 6: Priority Actions ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Priority Actions (Next 90 Days)"

    content = slide.placeholders[1]
    content.text = (
        "1. Optimize high-cost services\n"
        "2. Implement Savings Plans / Reserved Instances\n"
        "3. Remove idle resources\n"
        "4. Establish FinOps governance"
    )

    # ---------- Slide 7: Investment & ROI ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Investment & ROI Outlook"

    content = slide.placeholders[1]
    content.text = (
        f"Projected Annual Spend: ₹{annual_spend:,.0f}\n"
        f"Potential Annual Savings: ₹{annual_savings:,.0f}\n"
        "ROI Timeline: < 12 months\n"
        "Savings can fund modernization initiatives"
    )

    # ---------- Slide 8: Transformation Roadmap ----------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Transformation Roadmap"

    content = slide.placeholders[1]
    content.text = (
        "Phase 1: Stabilize & Optimize (0–6 months)\n"
        "Phase 2: Modernize & Automate (6–18 months)\n"
        "Phase 3: Innovate & Transform (18+ months)"
    )

    output_path = "Partner_Level_Board_Pack.pptx"
    prs.save(output_path)

    return output_path


print("PowerPoint report generator loaded")
