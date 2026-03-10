from pptx import Presentation
from pptx.util import Inches, Pt


def generate_mckinsey_deck(
    client,
    monthly_spend,
    savings_monthly,
    maturity_score,
    readiness_score,
    top_service
):

    annual_spend = monthly_spend * 12
    annual_savings = savings_monthly * 12

    prs = Presentation()

    # ===== Slide 1 — Title =====
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Cloud Cost & Transformation Review"
    slide.placeholders[1].text = f"{client}"

    # ===== Slide 2 — Executive Takeaway =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Takeaway"

    takeaway = (
        f"Cloud spend is concentrated in {top_service}, "
        f"creating a ₹{annual_savings:,.0f} annual optimization opportunity "
        "achievable within 12 months."
    )

    slide.placeholders[1].text = takeaway

    # ===== Slide 3 — Current State =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Current Cloud State"

    slide.placeholders[1].text = (
        f"Monthly Spend: ₹{monthly_spend:,.0f}\n"
        f"Cloud Maturity: {maturity_score}/100\n"
        f"Transformation Readiness: {readiness_score}/100"
    )

    # ===== Slide 4 — Key Issues =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Cost & Governance Issues"

    slide.placeholders[1].text = (
        "• High concentration in limited services\n"
        "• Potential overprovisioning\n"
        "• Governance improvements required"
    )

    # ===== Slide 5 — Recommended Actions =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recommended Actions"

    slide.placeholders[1].text = (
        "1. Optimize high-cost services\n"
        "2. Implement Savings Plans\n"
        "3. Remove idle resources"
    )

    # ===== Slide 6 — Financial Impact =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Financial Impact"

    slide.placeholders[1].text = (
        f"Projected Annual Spend: ₹{annual_spend:,.0f}\n"
        f"Potential Annual Savings: ₹{annual_savings:,.0f}\n"
        "ROI Timeline: < 12 months"
    )

    # ===== Slide 7 — Transformation Roadmap =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Transformation Roadmap"

    slide.placeholders[1].text = (
        "Phase 1: Stabilize & Optimize\n"
        "Phase 2: Modernize & Automate\n"
        "Phase 3: Innovate & Transform"
    )

    output = "McKinsey_Style_Cloud_Deck.pptx"
    prs.save(output)

    return output
